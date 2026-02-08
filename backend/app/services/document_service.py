"""Service for processing documents (PDF, DOCX, PPTX) and extracting text."""

import os
import tempfile
from pathlib import Path


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file."""
    try:
        from pypdf import PdfReader
        
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"Failed to extract text from PDF: {str(e)}")


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a DOCX file."""
    try:
        from docx import Document
        
        doc = Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"Failed to extract text from DOCX: {str(e)}")


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"Failed to extract text from PPTX: {str(e)}")


def extract_text_from_document(file_path: str, file_ext: str) -> str:
    """
    Extract text from a document file.
    
    Args:
        file_path: Path to the document file
        file_ext: File extension (pdf, docx, pptx)
    
    Returns:
        Extracted text content
    
    Raises:
        Exception: If file format is not supported or extraction fails
    """
    file_ext = file_ext.lower()
    
    if file_ext == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_ext in ['docx', 'doc']:
        return extract_text_from_docx(file_path)
    elif file_ext in ['pptx', 'ppt']:
        return extract_text_from_pptx(file_path)
    else:
        raise Exception(f"Unsupported file format: .{file_ext}")
